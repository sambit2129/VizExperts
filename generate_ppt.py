from pptx import Presentation
from pptx.util import Inches

# Load transcription
with open("script.txt", "r") as file:
    transcription = file.read()

# Create a presentation
prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
slide.shapes.title.text = "Industrial Equipment Overview"
slide.placeholders[1].text = "Generated using AI/ML pipeline"

# Slide 2: Script
slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
slide.shapes.title.text = "Generated Script"
slide.placeholders[1].text = transcription

# Slide 3: Image
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
slide.shapes.title.text = "Equipment Image"
slide.shapes.add_picture("frame_30.jpg", Inches(1), Inches(2), Inches(4), Inches(3))

# Save presentation
prs.save("output_presentation.pptx")
print("PowerPoint saved as output_presentation.pptx")
